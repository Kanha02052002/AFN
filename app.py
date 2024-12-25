from flask import Flask, request, render_template, send_from_directory, redirect, url_for
import os
from werkzeug.utils import secure_filename
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from fpdf import FPDF
import PyPDF2
import fitz  # PyMuPDF

# Configurations
app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['ALLOWED_EXTENSIONS'] = {'pdf', 'docx', 'pptx', 'txt'}
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # Max file size 16MB

# Check if file extension is allowed
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

# Delete all files in the upload folder
def delete_uploaded_files():
    for filename in os.listdir(app.config['UPLOAD_FOLDER']):
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        try:
            if os.path.isfile(file_path):
                os.unlink(file_path)
        except Exception as e:
            print(f"Error deleting file {file_path}: {e}")

# Convert PDF to DOCX (including images)
def convert_pdf_to_docx(input_file, original_filename):
    output_filename = os.path.join(app.config['UPLOAD_FOLDER'], f"{os.path.splitext(original_filename)[0]}_converted.docx")
    doc = Document()
    pdf_document = fitz.open(input_file)

    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        text = page.get_text("text")
        doc.add_paragraph(text)

        # Extract images
        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = pdf_document.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            image_filename = os.path.join(app.config['UPLOAD_FOLDER'], f"image_{page_num}_{img_index}.{image_ext}")
            with open(image_filename, "wb") as image_file:
                image_file.write(image_bytes)
            doc.add_picture(image_filename, width=Inches(5.0))

    doc.save(output_filename)
    return f"{os.path.splitext(original_filename)[0]}_converted.docx"

# Convert PDF to PPTX (simple conversion, may lose formatting)
def convert_pdf_to_pptx(input_file, original_filename):
    output_filename = os.path.join(app.config['UPLOAD_FOLDER'], f"{os.path.splitext(original_filename)[0]}_converted.pptx")
    with open(input_file, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        presentation = Presentation()
        for page in reader.pages:
            slide = presentation.slides.add_slide(presentation.slide_layouts[0])
            textbox = slide.shapes.add_textbox(left=100, top=100, width=600, height=400)
            textbox.text = page.extract_text()
        presentation.save(output_filename)
    return f"{os.path.splitext(original_filename)[0]}_converted.pptx"

# Convert DOCX to PDF
def convert_docx_to_pdf(input_file, original_filename):
    doc = Document(input_file)
    pdf_filename = os.path.join(app.config['UPLOAD_FOLDER'], f"{os.path.splitext(original_filename)[0]}_converted.pdf")
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    for paragraph in doc.paragraphs:
        pdf.set_font('Arial', size=12)
        pdf.multi_cell(0, 10, paragraph.text)
    
    pdf.output(pdf_filename)
    return f"{os.path.splitext(original_filename)[0]}_converted.pdf"

# Convert DOCX to PPTX
def convert_docx_to_pptx(input_file, original_filename):
    doc = Document(input_file)
    pptx_filename = os.path.join(app.config['UPLOAD_FOLDER'], f"{os.path.splitext(original_filename)[0]}_converted.pptx")
    presentation = Presentation()
    for paragraph in doc.paragraphs:
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        textbox = slide.shapes.add_textbox(left=100, top=100, width=600, height=400)
        textbox.text = paragraph.text
    presentation.save(pptx_filename)
    return f"{os.path.splitext(original_filename)[0]}_converted.pptx"

# Convert PPTX to PDF
def convert_pptx_to_pdf(input_file, original_filename):
    presentation = Presentation(input_file)
    pdf_filename = os.path.join(app.config['UPLOAD_FOLDER'], f"{os.path.splitext(original_filename)[0]}_converted.pdf")
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)

    for slide in presentation.slides:
        pdf.add_page()
        pdf.set_font('Arial', size=12)
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                pdf.multi_cell(0, 10, shape.text)
    
    pdf.output(pdf_filename)
    return f"{os.path.splitext(original_filename)[0]}_converted.pdf"

# Convert PPTX to DOCX
def convert_pptx_to_docx(input_file, original_filename):
    presentation = Presentation(input_file)
    doc_filename = os.path.join(app.config['UPLOAD_FOLDER'], f"{os.path.splitext(original_filename)[0]}_converted.docx")
    doc = Document()
    
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                try:
                    doc.add_paragraph(shape.text)
                except UnicodeEncodeError as e:
                    print(f"Error encoding text: {e}")
                    doc.add_paragraph(shape.text.encode('utf-8', errors='ignore').decode('utf-8'))

    doc.save(doc_filename)
    return f"{os.path.splitext(original_filename)[0]}_converted.docx"

# Convert TXT to DOCX
def convert_txt_to_docx(input_file, original_filename):
    doc_filename = os.path.join(app.config['UPLOAD_FOLDER'], f"{os.path.splitext(original_filename)[0]}_converted.docx")
    doc = Document()
    with open(input_file, 'r', encoding='utf-8') as file:
        for line in file:
            doc.add_paragraph(line)
    doc.save(doc_filename)
    return f"{os.path.splitext(original_filename)[0]}_converted.docx"

# Convert TXT to PDF
def convert_txt_to_pdf(input_file, original_filename):
    pdf_filename = os.path.join(app.config['UPLOAD_FOLDER'], f"{os.path.splitext(original_filename)[0]}_converted.pdf")
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font('Arial', size=12)
    with open(input_file, 'r', encoding='utf-8') as file:
        for line in file:
            pdf.multi_cell(0, 10, line)
    pdf.output(pdf_filename)
    return f"{os.path.splitext(original_filename)[0]}_converted.pdf"

# Convert TXT to PPTX
def convert_txt_to_pptx(input_file, original_filename):
    pptx_filename = os.path.join(app.config['UPLOAD_FOLDER'], f"{os.path.splitext(original_filename)[0]}_converted.pptx")
    presentation = Presentation()
    with open(input_file, 'r', encoding='utf-8') as file:
        for line in file:
            slide = presentation.slides.add_slide(presentation.slide_layouts[0])
            textbox = slide.shapes.add_textbox(left=100, top=100, width=600, height=400)
            textbox.text = line
    presentation.save(pptx_filename)
    return f"{os.path.splitext(original_filename)[0]}_converted.pptx"

@app.route('/', methods=['GET', 'POST'])
def index():
    delete_uploaded_files()  # Delete all uploaded files on page refresh
    if request.method == 'POST':
        if 'file' not in request.files:
            return redirect(request.url)
        file = request.files['file']
        if file.filename == '' or not allowed_file(file.filename):
            return redirect(request.url)
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        return render_template('index.html', uploaded=True, filename=filename, file_path=file_path)
    return render_template('index.html')

@app.route('/convert', methods=['POST'])
def convert():
    file_path = request.form['file_path']
    original_filename = os.path.basename(file_path)
    formats = request.form.getlist('formats')
    converted_files = []

    if 'pdf' in formats and not file_path.endswith('.pdf'):
        if file_path.endswith('.docx'):
            converted_files.append(convert_docx_to_pdf(file_path, original_filename))
        elif file_path.endswith('.pptx'):
            converted_files.append(convert_pptx_to_pdf(file_path, original_filename))
        elif file_path.endswith('.txt'):
            converted_files.append(convert_txt_to_pdf(file_path, original_filename))
    if 'docx' in formats and not file_path.endswith('.docx'):
        if file_path.endswith('.pdf'):
            converted_files.append(convert_pdf_to_docx(file_path, original_filename))
        elif file_path.endswith('.pptx'):
            converted_files.append(convert_pptx_to_docx(file_path, original_filename))
        elif file_path.endswith('.txt'):
            converted_files.append(convert_txt_to_docx(file_path, original_filename))
    if 'pptx' in formats and not file_path.endswith('.pptx'):
        if file_path.endswith('.pdf'):
            converted_files.append(convert_pdf_to_pptx(file_path, original_filename))
        elif file_path.endswith('.docx'):
            converted_files.append(convert_docx_to_pptx(file_path, original_filename))
        elif file_path.endswith('.txt'):
            converted_files.append(convert_txt_to_pptx(file_path, original_filename))

    return render_template('index.html', converted_files=converted_files)

@app.route('/download/<filename>')
def download(filename):
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename)

@app.route('/restart', methods=['POST'])
def restart():
    delete_uploaded_files()
    return redirect(url_for('index'))

if __name__ == '__main__':
    if not os.path.exists(app.config['UPLOAD_FOLDER']):
        os.makedirs(app.config['UPLOAD_FOLDER'])
    app.run(debug=True)